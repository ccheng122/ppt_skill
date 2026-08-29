"""Thin python-pptx wrappers and a shared palette for building slide decks.

Import these into a build script so every deck looks like one system and you do
not re-author text boxes each time. All positions are in inches on a 13.333 x 7.5
widescreen slide. See example_build.py for usage.
"""
import colorsys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- neutrals and semantic colors (a generic, reusable base) ----
INK      = RGBColor(0x15, 0x20, 0x2A)   # primary text
INK2     = RGBColor(0x48, 0x59, 0x65)   # secondary text
FAINT    = RGBColor(0x7B, 0x8B, 0x96)   # captions, footnotes
CARD     = RGBColor(0xF4, 0xF7, 0xF7)   # panel fill
BORDER   = RGBColor(0xDB, 0xE3, 0xE7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GOOD     = RGBColor(0x2F, 0x7A, 0x54)
GOOD_S   = RGBColor(0xDB, 0xED, 0xE2)
WARN     = RGBColor(0x9C, 0x6A, 0x15)
WARN_S   = RGBColor(0xF1, 0xE7, 0xCE)
GREY     = RGBColor(0x60, 0x6C, 0x74)

# ---- accent options ----
# Each option is an accent, a darker accent-ink for text on a soft band, a soft fill
# for callout bands, and a hex for matching charts (see chart_helpers). Pick the one
# that fits the subject or the brand the deck mirrors, or add your own. None is a
# default house style; choose deliberately per deck.
def _rgb(h): return RGBColor((h >> 16) & 255, (h >> 8) & 255, h & 255)
PALETTES = {
    "teal":    dict(accent=_rgb(0x17706B), ink=_rgb(0x0F4E4A), soft=_rgb(0xD6E8E5), hex="#17706B"),
    "indigo":  dict(accent=_rgb(0x2B4F8A), ink=_rgb(0x1B335C), soft=_rgb(0xDDE5F3), hex="#2B4F8A"),
    "plum":    dict(accent=_rgb(0x7A3B6B), ink=_rgb(0x54284A), soft=_rgb(0xEEDDE9), hex="#7A3B6B"),
    "forest":  dict(accent=_rgb(0x2F6E4F), ink=_rgb(0x204C37), soft=_rgb(0xDBEBE1), hex="#2F6E4F"),
    "rust":    dict(accent=_rgb(0xB0502E), ink=_rgb(0x7A371E), soft=_rgb(0xF3E1D7), hex="#B0502E"),
    "slate":   dict(accent=_rgb(0x3A4A57), ink=_rgb(0x26323C), soft=_rgb(0xE2E8EC), hex="#3A4A57"),
}

# Neutral-leaning vs. expressive, for when there is no brand to match (see
# use_palette / palette_from_hex docstrings below). Reach for a neutral one by
# default; reach for plum or rust only when the subject calls for more warmth.
NEUTRAL_PALETTES = ("slate", "teal", "forest", "indigo")
EXPRESSIVE_PALETTES = ("plum", "rust")

# Current accent. Call use_palette(...) in your build script to set it before laying
# out slides. Defaults to a neutral slate so nothing colorful is imposed by accident.
ACCENT, ACCENT_INK, ACCENT_SOFT = PALETTES["slate"]["accent"], PALETTES["slate"]["ink"], PALETTES["slate"]["soft"]


def use_palette(name):
    """Set the active accent from PALETTES (by name) or a custom dict {accent,ink,soft}."""
    global ACCENT, ACCENT_INK, ACCENT_SOFT
    p = PALETTES[name] if isinstance(name, str) else name
    ACCENT, ACCENT_INK, ACCENT_SOFT = p["accent"], p["ink"], p["soft"]
    return p


def _hex_to_hls(hex_color):
    hex_color = hex_color.lstrip("#")
    r, g, b = (int(hex_color[i:i + 2], 16) / 255 for i in (0, 2, 4))
    return colorsys.rgb_to_hls(r, g, b)


def _hls_to_hex(h, l, s):
    r, g, b = colorsys.hls_to_rgb(h, l, s)
    return "#{:02X}{:02X}{:02X}".format(round(r * 255), round(g * 255), round(b * 255))


def palette_from_hex(brand_hex):
    """Derive a deck-ready {accent, ink, soft, hex} palette from one brand color.

    Use this when the deck is for a specific company: seed it with their
    primary brand color and it keeps that hue, so the deck reads as theirs
    rather than getting one of the six house options pasted on. (If the brand
    has several colors, pick the dominant one; the deck stays one accent per
    its own restraint rule, it does not try to reproduce a whole brand system.)

    It only corrects the two failure modes that make a raw brand hex read
    badly at deck scale: a color too light, or too washed out/low-saturation,
    to hold up as small text or a thin chart line on white. A well-behaved
    brand color comes back close to unchanged; a pastel or muted one gets
    pulled just far enough to stay legible.

    For a generic or internal deck with no brand to match, skip this and call
    use_palette() with one of NEUTRAL_PALETTES instead.
    """
    h, l, s = _hex_to_hls(brand_hex)

    a_l = min(l, 0.55)
    a_s = min(max(s, 0.35), 0.85)
    accent_hex = _hls_to_hex(h, a_l, a_s)

    ink_l = max(a_l * 0.5, 0.14)
    ink_s = min(a_s + 0.05, 0.85)
    ink_hex = _hls_to_hex(h, ink_l, ink_s)

    soft_hex = _hls_to_hex(h, 0.91, min(a_s * 0.4, 0.32))

    return dict(
        accent=RGBColor.from_string(accent_hex.lstrip("#")),
        ink=RGBColor.from_string(ink_hex.lstrip("#")),
        soft=RGBColor.from_string(soft_hex.lstrip("#")),
        hex=accent_hex,
    )

FONT = "Arial"  # portable across PowerPoint and Google Slides


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill, line=None, round_=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = 0.05
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def R(t, size, color=INK, bold=False, spacing=None, caps=False):
    """A styled text run tuple. spacing is letter spacing in points; caps upper-cases."""
    return (t, size, color, bold, spacing, caps)


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        sp_after=4, line_sp=1.1):
    """runs is a list of paragraphs; each paragraph is a list of R(...) tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        p.line_spacing = line_sp
        for (t, size, color, bold, spacing, caps) in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(size)
            f.color.rgb = color
            f.bold = bold
            f.name = FONT
            if spacing is not None:
                r._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
            if caps:
                r._r.get_or_add_rPr().set('cap', 'all')
    return tb


def pic(s, path, x, y, w):
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def add_soft_shadow(shape, blur_pt=9, dist_pt=1.5, alpha_pct=20, dir_deg=90):
    """Attach a soft outer drop shadow to any autoshape (spec pulled from a
    real reference deck's chart-card treatment: 9pt blur, 1.5pt straight-down
    offset, black at 20% alpha). python-pptx has no high-level shadow API for
    a custom effect, so this writes the a:effectLst/a:outerShdw XML directly;
    call it right after creating the shape, before adding anything on top."""
    spPr = shape._element.spPr
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur_pt * 12700)),
        'dist': str(int(dist_pt * 12700)),
        'dir': str(int(dir_deg * 60000)),
        'rotWithShape': '0',
        'algn': 'bl',
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_pct * 1000)})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def chart_card(s, x, y, w, h, pad=0.15, corner=0.06):
    """White rounded-rect frame + soft drop shadow behind a chart image, so
    the chart reads as a card rather than floating loose on the slide (a
    pattern pulled from a real reference deck, not invented). Call this
    BEFORE placing the picture with pic(x, y, w) at the same x/y/w/h, so the
    picture renders on top of the card. h should be the chart's own display
    height (w / the source image's aspect ratio), not counting pad.

    A callout/band placed directly under one or more chart cards should span
    the cards' combined outer footprint (x - pad to the last card's x + w +
    pad), not the slide's full content width, or the edges visibly mismatch.
    """
    shp = rect(s, x - pad, y - pad, w + 2 * pad, h + 2 * pad, WHITE, round_=True)
    try:
        shp.adjustments[0] = corner
    except Exception:
        pass
    shp.line.fill.background()
    add_soft_shadow(shp)
    return shp


def eyebrow(s, text, kicker=None, x0=0.55):
    """Section label at the top of a slide, with an optional right-aligned kicker.

    x0 is the content left margin; pass it through when a left section_bar()
    (or anything else) has shifted content off the default 0.55 margin.
    """
    rect(s, x0, 0.5, 0.34, 0.05, ACCENT, round_=False)
    txt(s, x0 + 0.4, 0.4, 8.55 - (x0 - 0.55), 0.35, [[R(text, 11.5, ACCENT, True, 1.4, True)]])


def section_bar(s, label, bar_w=0.42, slide_w=13.333, slide_h=7.5, font_size=14):
    """Full-height color rail on the left edge, rotated text labeling which
    part of the story this slide is (a real pattern from a reference deck's
    per-section spine: a rotated rectangle, one color per section, walking
    the deck as you go). Optional, not a default: offer it the way
    divider()/agenda() are offered, don't add it unprompted.

    A single narrative deck (most decks here) has no real sections to color
    differently, so default to ONE accent color throughout with the label
    text changing per slide (e.g. eyebrow's own text) rather than inventing
    arbitrary per-slide colors that carry no real meaning. Reserve distinct
    colors per call for a deck that actually has distinct named sections.

    Content on the slide (eyebrow, footer, headline, everything) needs its
    left margin shifted right to clear the bar; a margin of bar_w + ~0.55in
    from the slide edge reads clean. Pass x0=<that margin> to eyebrow() and
    footer(), and use the same left x for every other element on the slide.

    Mechanics: this is a plain rectangle rotated -90 degrees about its own
    center (rot="-5400000" in the underlying XML), which is how a reference
    deck built the same effect. That rotation maps the unrotated box's LEFT
    edge to the bar's BOTTOM and its RIGHT edge to the bar's TOP (confirmed
    against that reference deck, which right-aligns its own rail text the
    same way). So right-aligning the paragraph here, plus a right-side inset,
    is what actually reads as "top-aligned" once the whole shape rotates;
    center alignment puts the label at the vertical middle of the slide,
    which reads as floating rather than a section tag.
    """
    cx, cy = bar_w / 2, slide_h / 2
    left, top = cx - slide_h / 2, cy - bar_w / 2
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                              Inches(slide_h), Inches(bar_w))
    shp.rotation = -90
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = Inches(0.4)  # inset from the box's right edge = gap from the bar's visual top
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = label
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    return shp
    if kicker:
        txt(s, 9.0, 0.38, 3.8, 0.35, [[R(kicker, 11, FAINT, False, 0.3)]], align=PP_ALIGN.RIGHT)


def band(s, y, runs, x=0.55, w=12.23, h=0.92, fill=None):
    """Full-width highlight bar for a verdict or callout. runs is one paragraph.

    Defaults to the active ACCENT_SOFT, resolved at call time so use_palette() applies.
    """
    rect(s, x, y, w, h, fill if fill is not None else ACCENT_SOFT)
    txt(s, x + 0.3, y, w - 0.6, h, [runs], anchor=MSO_ANCHOR.MIDDLE, line_sp=1.1)


def bullets(s, x, y, w, h, items, size=14, color=INK, marker="•",
            sp_after=8, line_sp=1.2, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """A flat list of bullet points, one idea per line, not a paragraph.

    items is a list of strings. A full sentence is fine; a bullet that
    splices two ideas together with a semicolon, or runs past ~2 lines,
    belongs back in prose or split into two bullets. No raw internal
    notation (formulas, code-shaped values) as a bullet, translate it to
    plain language first. Keep sibling bullets parallel in phrasing.
    """
    runs = [[R(f"{marker}  ", size, ACCENT, True), R(item, size, color)] for item in items]
    return txt(s, x, y, w, h, runs, align=align, anchor=anchor, sp_after=sp_after, line_sp=line_sp)


def divider(s, title, number=None, kicker=None):
    """Section-break slide: an optional numeral, a big title, roughly vertically
    centered. Use between major sections in a longer, multi-section deck, not
    in the default 3-5 slide shape (that would just be padding).

    number is a short label like "01" or "I"; omit it for a plain title card.
    """
    cy = 3.15
    if number:
        txt(s, 0.9, cy, 2.0, 1.1, [[R(number, 52, ACCENT, True)]])
        txt(s, 2.75, cy + 0.2, 9.6, 1.1, [[R(title, 38, INK, True)]])
    else:
        txt(s, 0.9, cy + 0.2, 11.5, 1.1, [[R(title, 38, INK, True)]])
    if kicker:
        kx = 2.75 if number else 0.9
        txt(s, kx, cy + 1.05, 9.6, 0.4, [[R(kicker, 13, FAINT)]])


def table(s, x, y, w, h, headers, rows, col_widths=None):
    """A simple styled comparison table, for findings that are genuinely a
    side-by-side comparison rather than a chart or prose.

    headers is a list of column labels. rows is a list of row-lists, each
    the same length as headers. col_widths, if given, is a list of relative
    weights (e.g. [2, 1, 1] to make the first column twice as wide); default
    is equal columns. Header row fills with the active ACCENT; body rows
    alternate white and CARD for scan-ability.
    """
    # A thin grey outline behind the table: without it, white body rows have
    # no edge and blend straight into the slide's white background (caught in
    # real use: a lone table on an otherwise-white slide read as floating).
    rect(s, x - 0.02, y - 0.02, w + 0.04, h + 0.04, WHITE, line=BORDER, round_=False)

    n_cols = len(headers)
    gshape = s.shapes.add_table(len(rows) + 1, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gshape.table
    tbl.first_row = False
    tbl.horz_banding = False

    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(w * cw / total)

    def _cell(i, j, value, fill, color, bold):
        cell = tbl.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(value)
        r.font.size = Pt(11)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT

    for j, head in enumerate(headers):
        _cell(0, j, head, ACCENT, WHITE, True)
    for i, row in enumerate(rows, start=1):
        fill = WHITE if i % 2 else CARD
        for j, value in enumerate(row):
            _cell(i, j, value, fill, INK, False)

    return gshape


def stat_group(s, x, y, stats, w=2.6, step=1.3, gap=0.4, direction="vertical",
               value_size=30, label_size=11):
    """A group of 2-3 big stats, value over a one-line label, meant to sit
    beside a chart placed with pic() (see example_build.py). stats is a
    list of (value, label) tuples, e.g. [("7.4%", "Of requests result in a
    support call")].

    direction="vertical" (default) stacks them top to bottom, step apart;
    direction="horizontal" lays them left to right, each w wide with gap
    between, e.g. across the top of a slide instead of beside a chart. No
    icons in this version, the number is the anchor; add an image per stat
    yourself with pic() if a deck genuinely needs one.
    """
    for i, (value, label) in enumerate(stats):
        if direction == "horizontal":
            bx, by = x + i * (w + gap), y
        else:
            bx, by = x, y + i * step
        txt(s, bx, by, w, 0.6, [[R(value, value_size, ACCENT, True)]])
        txt(s, bx, by + 0.56, w, 0.6, [[R(label, label_size, INK2)]], line_sp=1.15)


def two_up(s, x, y, w, h, items, gap=0.5, label_size=11):
    """Two chart images side by side, each with a small caption above it,
    for a direct visual comparison (e.g. before/after, two segments, two
    markets). items is a list of exactly 2 (image_path, label) tuples.
    Speculative, still unconfirmed against a real deck, watch whether the
    fixed even-width split actually fits real chart aspect ratios in
    practice before leaning on it.
    """
    col_w = (w - gap) / 2
    for i, (img_path, label) in enumerate(items):
        cx = x + i * (col_w + gap)
        txt(s, cx, y, col_w, 0.35, [[R(label, label_size, FAINT, True, 0.4, True)]])
        pic(s, img_path, cx, y + 0.4, col_w)


def footer(s, text, page=None, x0=0.55):
    """Small bottom-of-slide caption: a source, date, or confidentiality
    line, with an optional page number right-aligned at the same baseline.

    x0 is the content left margin; pass it through alongside eyebrow(x0=...).
    """
    txt(s, x0, 7.05, 10.6 - (x0 - 0.55), 0.3, [[R(text, 9, FAINT)]])
    if page is not None:
        txt(s, 11.3, 7.05, 1.48, 0.3, [[R(str(page), 9, FAINT)]], align=PP_ALIGN.RIGHT)


def agenda(s, items, title="Agenda"):
    """Numbered section list for the front of a longer, multi-section deck.
    items is a list of section-title strings, numbered in order.
    """
    eyebrow(s, title)
    y = 1.5
    for i, item in enumerate(items, start=1):
        txt(s, 0.55, y, 0.7, 0.5, [[R(f"{i:02d}", 16, ACCENT, True)]])
        txt(s, 1.35, y, 10.9, 0.5, [[R(item, 18, INK)]])
        y += 0.62
